"""Generate the VeriLens hackathon presentation (.pptx).

Run:  .venv\\Scripts\\python.exe make_ppt.py
Output: VeriLens_Presentation.pptx
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------------------
# Theme (matches the VeriLens UI)
# ---------------------------------------------------------------------------

LAVENDER_BG = RGBColor(0xF8, 0xF7, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1E, 0x1B, 0x3A)      # dark navy/purple text
PURPLE = RGBColor(0x6D, 0x28, 0xD9)
PURPLE_DARK = RGBColor(0x4C, 0x1D, 0x95)
PINK = RGBColor(0xEC, 0x48, 0x99)
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LIGHT_CARD = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE5, 0xE7, 0xEB)

FONT = "Calibri"
FONT_BOLD = "Calibri"

EMU_IN = 914400


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def blank_slide(prs: Presentation):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # lavender background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = LAVENDER_BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    color: RGBColor,
    radius: bool = False,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if radius:
        try:
            shape.adjustments[0] = 0.08
        except Exception:
            pass
    return shape


def add_text(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = DARK,
    align: int = PP_ALIGN.LEFT,
    anchor: int = MSO_ANCHOR.TOP,
    italic: bool = False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    return box


def add_bullets(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    items: list[tuple[str, str]],  # (bullet, color) pairs; color "" -> dark
    *,
    size: int = 17,
    gap: int = 10,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    first = True
    for bullet, color in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        run = p.add_run()
        run.text = f"▪  {bullet}"
        f = run.font
        f.name = FONT
        f.size = Pt(size)
        f.color.rgb = RGBColor.from_string(color) if color else DARK
    return box


def header(slide, kicker: str, title: str, number: str):
    """Standard content-slide header: purple accent bar + kicker + title."""
    add_rect(slide, 0.55, 0.55, 0.12, 0.95, color=PURPLE)
    add_text(slide, 0.85, 0.55, 9.0, 0.4, kicker.upper(), size=13, bold=True, color=PURPLE)
    add_text(slide, 0.85, 0.88, 10.5, 0.7, title, size=30, bold=True, color=DARK)
    add_text(slide, 12.2, 0.6, 0.8, 0.4, number, size=14, bold=True, color=MUTED, align=PP_ALIGN.RIGHT)
    add_rect(slide, 0.55, 1.78, 12.23, 0.02, color=CARD_BORDER)


def footer(slide):
    add_text(
        slide,
        0.55,
        7.08,
        12.23,
        0.3,
        "VeriLens · AI-Powered Fake News Detector · Algorithm Avengers",
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_screenshot(slide, path: str, x: float, y: float, max_w: float, max_h: float):
    pic = slide.shapes.add_picture(path, Inches(x), Inches(y))
    scale = min(Inches(max_w) / pic.width, Inches(max_h) / pic.height)
    pic.width = Emu(int(pic.width * scale))
    pic.height = Emu(int(pic.height * scale))
    # recenter horizontally
    slide_w = Inches(13.333)
    pic.left = Emu(int((slide_w - pic.width) / 2))
    return pic


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

prs = new_presentation()
SW = 13.333

# ---- Slide 1: Title -------------------------------------------------------
s = blank_slide(prs)
add_rect(s, 0, 0, SW, 7.5, color=PURPLE_DARK)
add_rect(s, 0, 7.28, SW, 0.22, color=PINK)
add_text(s, 1.0, 1.7, 11.33, 0.5, "AI-POWERED FACT CHECKING", size=15, bold=True,
         color=RGBColor(0xC4, 0xB5, 0xFD), align=PP_ALIGN.CENTER)
add_text(s, 1.0, 2.3, 11.33, 1.1, "VeriLens", size=66, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, 1.0, 3.55, 11.33, 0.6, "Verify Before You Believe.", size=26,
         color=RGBColor(0xF9, 0xA8, 0xD4), align=PP_ALIGN.CENTER, italic=True)
add_text(s, 1.0, 4.45, 11.33, 0.5, "An AI-powered fake news detector: retrieve real evidence, "
         "verify the claim, explain the verdict.", size=17, color=RGBColor(0xDD, 0xD6, 0xFE),
         align=PP_ALIGN.CENTER)
add_rect(s, 4.87, 5.45, 3.6, 0.035, color=PINK)
add_text(s, 1.0, 5.75, 11.33, 0.5, "Shree Deore   ·   Krish Dalvi   ·   Vaishnavi Rathi",
         size=19, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 6.3, 11.33, 0.4, "Built with ❤️ by Algorithm Avengers", size=14,
         color=RGBColor(0xA7, 0x8B, 0xFA), align=PP_ALIGN.CENTER)

# ---- Slide 2: The Problem --------------------------------------------------
s = blank_slide(prs)
header(s, "Why this matters", "The Problem", "02")
add_bullets(s, 0.9, 2.15, 11.5, 3.2, [
    ("Misinformation spreads faster than verified facts — it reaches more people, sooner.", "1E1B3A"),
    ("Readers cannot manually verify every claim they see on social media.", "1E1B3A"),
    ("Viral fake news erodes public trust in real journalism and institutions.", "1E1B3A"),
    ("Most people lack a fast, trustworthy way to check a claim before sharing it.", "1E1B3A"),
], size=19, gap=16)
add_rect(s, 0.9, 5.6, 11.53, 1.0, color=RGBColor(0xED, 0xE9, 0xFF), radius=True)
add_text(s, 1.25, 5.85, 10.8, 0.5,
         "“Falsehoods spread significantly faster than the truth online.”",
         size=17, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
add_text(s, 1.25, 6.3, 10.8, 0.3, "— MIT Media Lab, “The Spread of True and False News Online”",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)
footer(s)

# ---- Slide 3: The Solution -------------------------------------------------
s = blank_slide(prs)
header(s, "Our approach", "The Solution — RAG-Based Fact Checking", "03")
add_bullets(s, 0.9, 2.15, 11.5, 3.4, [
    ("VeriLens verifies a claim by retrieving real web evidence — it never invents facts.", "1E1B3A"),
    ("A RAG (Retrieval-Augmented Generation) pipeline extracts relevant passages and classifies them as supporting, contradicting or neutral.", "1E1B3A"),
    ("A deterministic verifier aggregates the evidence into a clear verdict (Supported / Contradicted / Mixed / Unverified).", "1E1B3A"),
    ("An LLM writes a short, evidence-grounded analysis — strictly from the retrieved material.", "1E1B3A"),
    ("If evidence is insufficient, VeriLens honestly says UNVERIFIED instead of guessing.", "6D28D9"),
], size=18, gap=13)
footer(s)

# ---- Slide 4: Technology Stack --------------------------------------------
s = blank_slide(prs)
header(s, "Tech we used", "Technology Stack", "04")

stack = [
    ("Backend / API", "Python 3 · FastAPI · Uvicorn", "REST API with automatic Swagger docs, Pydantic validation, CORS and error handling."),
    ("Retrieval", "Tavily / Serper web search", "Real-time evidence from the web with source authority ranking."),
    ("LLM", "OpenAI (GPT-4o-mini) / Groq", "Grounded explanations; deterministic rule-based fallback when offline."),
    ("Frontend", "HTML · CSS · Vanilla JS", "Fast, dependency-free chat & verification UI with dark mode."),
    ("Testing", "pytest (158 tests)", "Offline, mocked, deterministic — no API keys required."),
    ("Security", "Pydantic · rate limiting", "Input validation, safe errors, prompt-injection guards, .env secrets."),
]
y = 2.05
for title, sub, desc in stack:
    add_rect(s, 0.9, y, 11.53, 0.78, color=WHITE, radius=True)
    add_text(s, 1.2, y + 0.09, 3.4, 0.35, title, size=15, bold=True, color=PURPLE)
    add_text(s, 1.2, y + 0.4, 4.4, 0.3, sub, size=12.5, bold=True, color=DARK)
    add_text(s, 4.9, y + 0.12, 7.2, 0.6, desc, size=12.5, color=MUTED)
    y += 0.84
footer(s)

# ---- Slide 5: How It Works ------------------------------------------------
s = blank_slide(prs)
header(s, "Under the hood", "How It Works — The Pipeline", "05")

steps = [
    ("1", "User Claim", "User enters a news claim in VeriLens"),
    ("2", "Web Retrieval", "Search APIs fetch real articles"),
    ("3", "Evidence Extraction", "Relevant passages are selected"),
    ("4", "Classification", "Supporting / contradicting / neutral"),
    ("5", "Verification", "Verdict + confidence is computed"),
    ("6", "LLM Analysis", "Grounded explanation is generated"),
    ("7", "Result", "Verdict, evidence & sources shown"),
]
x = 0.55
for num, title, desc in steps:
    add_rect(s, x, 2.35, 1.62, 2.3, color=WHITE, radius=True)
    add_text(s, x + 0.1, 2.5, 1.42, 0.5, num, size=22, bold=True, color=PINK, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.08, 2.95, 1.46, 0.7, title, size=13.5, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.12, 3.55, 1.38, 1.0, desc, size=10, color=MUTED, align=PP_ALIGN.CENTER)
    if x + 1.62 < SW - 1.0:
        add_text(s, x + 1.6, 3.0, 0.25, 0.5, "›", size=20, bold=True, color=PURPLE,
                 align=PP_ALIGN.CENTER)
    x += 1.82

add_rect(s, 0.9, 5.15, 11.53, 1.15, color=RGBColor(0xED, 0xE9, 0xFF), radius=True)
add_text(s, 1.2, 5.35, 11.0, 0.75,
         "Evidence first. The LLM only sees retrieved evidence — never its own knowledge — "
         "and every citation is cross-checked against real sources.", size=15, bold=True,
         color=PURPLE, align=PP_ALIGN.CENTER)
footer(s)

# ---- Slides 6-9: Screenshots ----------------------------------------------
shots = [
    ("Landing Page", "screen_landing.png",
     "Clean hero: “Verify Before You Believe.” with a one-click claim input."),
    ("Entering a Claim", "screen_input.png",
     "Typing a claim — live character counter and the Verify Claim action."),
    ("Verification Result", "screen_result.png",
     "Verdict, confidence, claim recap, AI analysis and evidence cards with source + similarity."),
    ("Hero & Product Overview", "screen_hero.png",
     "Full landing experience with team footer (“Built with ❤️ by Algorithm Avengers”)."),
]
for i, (title, img, caption) in enumerate(shots, start=6):
    s = blank_slide(prs)
    header(s, "Product screens", f"Screenshot — {title}", f"0{i}")
    add_screenshot(s, f"ppt_assets/{img}", 0, 2.05, 11.6, 4.6)
    add_text(s, 0.9, 6.72, 11.53, 0.4, caption, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s)

# ---- Slide 10: Key Features -----------------------------------------------
s = blank_slide(prs)
header(s, "What you can do", "Key Features", "10")
features = [
    ("Verdict + Confidence", "Clear verdict badge with a confidence meter"),
    ("Evidence Cards", "Source and similarity shown for every passage"),
    ("AI Analysis", "Short, evidence-grounded explanation for each claim"),
    ("Chatbot", "Follow-up questions like “Why?” or “Show evidence”"),
    ("Dark Mode & History", "Conversation persists across reloads"),
    ("API & Docs", "Public REST API with Swagger docs (/docs)"),
]
x = 0.9
y = 2.15
for i, (title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    add_rect(s, x + col * 5.95, y + row * 1.55, 5.6, 1.35, color=WHITE, radius=True)
    add_text(s, x + col * 5.95 + 0.25, y + row * 1.55 + 0.15, 5.1, 0.4, title, size=16,
             bold=True, color=PURPLE)
    add_text(s, x + col * 5.95 + 0.25, y + row * 1.55 + 0.55, 5.1, 0.7, desc, size=12.5,
             color=MUTED)
footer(s)

# ---- Slide 11: Security ----------------------------------------------------
s = blank_slide(prs)
header(s, "Built safely", "Security", "11")
add_bullets(s, 0.9, 2.15, 11.5, 4.2, [
    ("Input validation — every claim is length-capped, trimmed and type-checked (HTTP 422 for bad input).", "1E1B3A"),
    ("Secrets stay in .env — never committed, never sent to the frontend, never in error messages.", "1E1B3A"),
    ("Prompt-injection protection — user claims and retrieved web text are treated as data, not instructions.", "1E1B3A"),
    ("Safe errors — internal stack traces and exception details never reach the user.", "1E1B3A"),
    ("Rate limiting on public endpoints (HTTP 429) to prevent accidental abuse.", "1E1B3A"),
    ("Configurable CORS with credentials disabled.", "1E1B3A"),
], size=17, gap=14)
footer(s)

# ---- Slide 12: Testing -----------------------------------------------------
s = blank_slide(prs)
header(s, "Quality", "Testing", "12")
add_rect(s, 0.9, 2.2, 3.4, 1.5, color=PURPLE, radius=True)
add_text(s, 0.9, 2.5, 3.4, 0.6, "158", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.9, 3.15, 3.4, 0.4, "automated tests · 0 failing", size=13, color=RGBColor(0xDD, 0xD6, 0xFE),
         align=PP_ALIGN.CENTER)
add_rect(s, 4.6, 2.2, 3.4, 1.5, color=PINK, radius=True)
add_text(s, 4.6, 2.5, 3.4, 0.6, "100%", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 4.6, 3.15, 3.4, 0.4, "offline & mocked — no API keys", size=13,
         color=RGBColor(0xFD, 0xE8, 0xF3), align=PP_ALIGN.CENTER)
add_rect(s, 8.3, 2.2, 4.13, 1.5, color=ORANGE, radius=True)
add_text(s, 8.3, 2.5, 4.13, 0.6, "pytest -q", size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 8.3, 3.15, 4.13, 0.4, "one command to verify everything", size=13,
         color=RGBColor(0xFF, 0xED, 0xD5), align=PP_ALIGN.CENTER)
add_bullets(s, 0.9, 4.15, 11.5, 2.5, [
    ("Retrieval, evidence extraction, classification, explainer and the full RAG pipeline.", "1E1B3A"),
    ("API layer: health, valid/invalid requests, safe errors, CORS and rate limiting.", "1E1B3A"),
    ("Chatbot, follow-ups, and security (prompt injection, secret exposure, unsafe content).", "1E1B3A"),
    ("Deterministic and repeatable — fast enough for every commit.", "6D28D9"),
], size=16, gap=11)
footer(s)

# ---- Slide 13: Demo Flow ---------------------------------------------------
s = blank_slide(prs)
header(s, "See it live", "Demo Flow", "13")
demo = [
    ("1", "Open VeriLens and type a claim — e.g. “India won the 2023 Cricket World Cup.”"),
    ("2", "Hit Verify Claim; watch retrieval, evidence extraction and verification run."),
    ("3", "Read the verdict, confidence and AI analysis."),
    ("4", "Open the evidence cards and inspect real sources."),
    ("5", "Ask a follow-up (“Why?”) and get an evidence-based answer."),
    ("6", "Try an uncertain claim → honest UNVERIFIED. Try malicious input → no secrets leak."),
]
y = 2.15
for num, text in demo:
    add_rect(s, 0.9, y, 0.55, 0.55, color=PURPLE, radius=True)
    add_text(s, 0.9, y + 0.08, 0.55, 0.4, num, size=16, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, 1.7, y + 0.06, 10.6, 0.55, text, size=16, color=DARK, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.72
footer(s)

# ---- Slide 14: Thank You ---------------------------------------------------
s = blank_slide(prs)
add_rect(s, 0, 0, SW, 7.5, color=PURPLE_DARK)
add_rect(s, 0, 0, SW, 0.22, color=PINK)
add_text(s, 1.0, 2.1, 11.33, 1.0, "Thank You!", size=56, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s, 1.0, 3.3, 11.33, 0.5, "Questions? Let's verify together.", size=22,
         color=RGBColor(0xF9, 0xA8, 0xD4), align=PP_ALIGN.CENTER, italic=True)
add_rect(s, 4.87, 4.2, 3.6, 0.035, color=PINK)
add_text(s, 1.0, 4.5, 11.33, 0.5, "Shree Deore   ·   Krish Dalvi   ·   Vaishnavi Rathi",
         size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 5.1, 11.33, 0.4, "Built with ❤️ by Algorithm Avengers  ·  VeriLens © 2026",
         size=14, color=RGBColor(0xA7, 0x8B, 0xFA), align=PP_ALIGN.CENTER)

prs.save("VeriLens_Presentation.pptx")
print("Saved VeriLens_Presentation.pptx with", len(prs.slides._sldIdLst), "slides")
