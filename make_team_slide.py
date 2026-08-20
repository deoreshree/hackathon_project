"""Generate the "Algorithm Avengers" team slide (PNG + PPTX).

Outputs:
  Team_Algorithm_Avengers.png   — 1920x1080 image
  Team_Algorithm_Avengers.pptx  — editable single slide
"""

from __future__ import annotations

import os

from PIL import Image, ImageDraw, ImageFont  # type: ignore

W, H = 1920, 1080

# Team
TEAM = [
    ("Shree Deore", "S", (0xA7, 0x8B, 0xFA)),   # purple
    ("Vaishnavi Rathi", "V", (0xF4, 0x72, 0xB6)),  # pink
    ("Krish Dalvi", "K", (0xFB, 0x92, 0x3C)),    # orange
]

TOP = (30, 27, 58)      # #1E1B3A
BOTTOM = (76, 29, 149)  # #4C1D95


def _font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        r"C:\Windows\Fonts\calibrib.ttf" if bold else r"C:\Windows\Fonts\calibri.ttf",
        r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
    ]
    for path in candidates:
        if os.path.exists(path):
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def _center_text(d: ImageDraw.ImageDraw, y: int, text: str, font, fill) -> int:
    w = d.textlength(text, font=font)
    d.text(((W - w) / 2, y), text, font=font, fill=fill)
    return int(w)


def build_png(path: str) -> None:
    # ---- Background gradient -----------------------------------------
    img = Image.new("RGB", (W, H))
    d = ImageDraw.Draw(img)
    for y in range(H):
        t = y / (H - 1)
        r = int(TOP[0] + (BOTTOM[0] - TOP[0]) * t)
        g = int(TOP[1] + (BOTTOM[1] - TOP[1]) * t)
        b = int(TOP[2] + (BOTTOM[2] - TOP[2]) * t)
        d.line([(0, y), (W, y)], fill=(r, g, b))

    # ---- Decorative glows ---------------------------------------------
    overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    od = ImageDraw.Draw(overlay)
    od.ellipse([-260, -280, 560, 540], fill=(236, 72, 153, 36))    # pink glow top-left
    od.ellipse([1400, -320, 2140, 420], fill=(167, 139, 250, 40))  # purple glow top-right
    od.ellipse([820, 700, 1220, 1100], fill=(251, 146, 60, 26))    # orange glow bottom
    img = Image.alpha_composite(img.convert("RGBA"), overlay)

    d = ImageDraw.Draw(img)

    # ---- Badge pill ----------------------------------------------------
    badge = "TEAM  ·  ALGORITHM AVENGERS"
    bf = _font(24, bold=True)
    bw = int(d.textlength(badge, font=bf)) + 70
    bx = (W - bw) / 2
    d.rounded_rectangle([bx, 92, bx + bw, 152], radius=30, fill=(255, 255, 255, 40))
    d.text(((W - d.textlength(badge, font=bf)) / 2, 106), badge, font=bf,
           fill=(255, 228, 230, 255))

    # ---- Title ----------------------------------------------------------
    tf = _font(104, bold=True)
    _center_text(d, 190, "Algorithm Avengers", tf, (255, 255, 255, 255))
    # accent underline
    d.rounded_rectangle([(W - 340) / 2, 350, (W + 340) / 2, 356], radius=3,
                        fill=(236, 72, 153, 255))

    # ---- Tagline ---------------------------------------------------------
    tg = _font(27)
    _center_text(d, 388, "Three members. One mission: verify before you believe.",
                 tg, (221, 214, 254, 255))

    # ---- Member cards ----------------------------------------------------
    card_w, card_h = 470, 320
    gap = 45
    x0 = (W - (3 * card_w + 2 * gap)) / 2
    y0 = 520

    for i, (name, initial, accent) in enumerate(TEAM):
        cx = x0 + i * (card_w + gap)
        card_overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
        cd = ImageDraw.Draw(card_overlay)
        cd.rounded_rectangle([cx, y0, cx + card_w, y0 + card_h], radius=28,
                             fill=(255, 255, 255, 26), outline=(255, 255, 255, 90), width=2)
        img = Image.alpha_composite(img, card_overlay)
        d = ImageDraw.Draw(img)

        # avatar circle with initial
        ax, ay, radius = cx + card_w / 2, y0 + 86, 66
        av_overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
        ad = ImageDraw.Draw(av_overlay)
        ad.ellipse([ax - radius, ay - radius, ax + radius, ay + radius], fill=accent + (255,))
        img = Image.alpha_composite(img, av_overlay)
        d = ImageDraw.Draw(img)
        init_font = _font(64, bold=True)
        iw = d.textlength(initial, font=init_font)
        d.text((ax - iw / 2, ay - 44), initial, font=init_font, fill=(255, 255, 255, 255))

        # name
        nf = _font(34, bold=True)
        nw = d.textlength(name, font=nf)
        d.text((cx + (card_w - nw) / 2, y0 + 186), name, font=nf, fill=(255, 255, 255, 255))

        # role-free label: core member
        rf = _font(20)
        label = "Core Member"
        lw = d.textlength(label, font=rf)
        d.text((cx + (card_w - lw) / 2, y0 + 246), label, font=rf, fill=(214, 211, 238, 255))

    # ---- Footer ----------------------------------------------------------
    ff = _font(18, bold=True)
    _center_text(d, 990, "VeriLens  ·  AI-Powered Fake News Detector  ·  Hackathon 2026",
                 ff, (167, 139, 250, 255))

    img.convert("RGB").save(path, "PNG")


# ---------------------------------------------------------------------------
# PPTX version
# ---------------------------------------------------------------------------

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402

DARK_P = RGBColor(0x1E, 0x1B, 0x3A)
PURPLE_P = RGBColor(0x4C, 0x1D, 0x95)
WHITE_P = RGBColor(0xFF, 0xFF, 0xFF)
PINK_P = RGBColor(0xEC, 0x48, 0x99)
LAVENDER_P = RGBColor(0xDD, 0xD6, 0xFE)
MUTED_P = RGBColor(0xA7, 0x8B, 0xFA)


def build_pptx(path: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])

    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_P
    bg.line.fill.background()
    bg.shadow.inherit = False

    def tb(x, y, w, h, text, *, size, bold=False, color=WHITE_P, align=PP_ALIGN.CENTER):
        box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.name = "Calibri"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        return box

    # Badge
    pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.9), Inches(0.7),
                              Inches(3.55), Inches(0.5))
    pill.fill.solid()
    pill.fill.fore_color.rgb = RGBColor(0x2E, 0x2A, 0x5C)
    pill.line.fill.background()
    pill.shadow.inherit = False
    tb(4.9, 0.76, 3.55, 0.35, "TEAM · ALGORITHM AVENGERS", size=13, bold=True,
       color=RGBColor(0xFF, 0xE4, 0xE6))

    # Title
    tb(1.0, 1.35, 11.33, 1.1, "Algorithm Avengers", size=54, bold=True)
    # accent line
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.0), Inches(2.6),
                              Inches(3.35), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = PINK_P
    line.line.fill.background()
    line.shadow.inherit = False

    # Tagline
    tb(1.0, 2.85, 11.33, 0.5, "Three members. One mission: verify before you believe.",
       size=18, color=LAVENDER_P)

    # Member cards
    names = ["Shree Deore", "Vaishnavi Rathi", "Krish Dalvi"]
    initials = ["S", "V", "K"]
    accents = [RGBColor(0xA7, 0x8B, 0xFA), RGBColor(0xF4, 0x72, 0xB6), RGBColor(0xFB, 0x92, 0x3C)]
    card_w = 3.7
    gap = 0.35
    x0 = (13.333 - (3 * card_w + 2 * gap)) / 2
    y0 = 3.75
    for i, name in enumerate(names):
        cx = x0 + i * (card_w + gap)
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(y0),
                                  Inches(card_w), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(0x2E, 0x2A, 0x5C)
        card.line.color.rgb = RGBColor(0x8B, 0x7E, 0xCF)
        card.line.width = Pt(1.25)
        card.shadow.inherit = False
        try:
            card.adjustments[0] = 0.06
        except Exception:
            pass

        # avatar
        av = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + card_w / 2 - 0.45),
                                Inches(y0 + 0.35), Inches(0.9), Inches(0.9))
        av.fill.solid()
        av.fill.fore_color.rgb = accents[i]
        av.line.fill.background()
        av.shadow.inherit = False
        tb(cx + card_w / 2 - 0.45, y0 + 0.47, 0.9, 0.7, initials[i], size=26, bold=True)

        tb(cx, y0 + 1.5, card_w, 0.5, name, size=19, bold=True)
        tb(cx, y0 + 2.0, card_w, 0.4, "Core Member", size=12, color=RGBColor(0xD6, 0xD3, 0xEE))

    # Footer
    tb(1.0, 6.85, 11.33, 0.4, "VeriLens  ·  AI-Powered Fake News Detector  ·  Hackathon 2026",
       size=11, color=MUTED_P)

    prs.save(path)


if __name__ == "__main__":
    build_png("Team_Algorithm_Avengers.png")
    build_pptx("Team_Algorithm_Avengers.pptx")
    print("Saved Team_Algorithm_Avengers.png and Team_Algorithm_Avengers.pptx")
