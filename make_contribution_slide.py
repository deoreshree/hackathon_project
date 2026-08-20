"""Generate the updated "Team Contribution" slide (3 members) as PPTX + PNG.

Outputs:
  Team_Contribution_Slide.pptx  — editable single slide
  Team_Contribution_Slide.png   — 1920x1080 image to paste into any deck
"""

from __future__ import annotations

import os
import textwrap
from pathlib import Path

# ---------------------------------------------------------------------------
# PPTX version
# ---------------------------------------------------------------------------

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x11, 0x18, 0x27)
PERIWINKLE = RGBColor(0xA5, 0xB4, 0xFC)
BLUE = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x4B, 0x55, 0x63)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE_DARK = RGBColor(0x1D, 0x4E, 0xD8)


def _tb(slide, x, y, w, h, text, *, size, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _card(slide, x, y, w, h, *, fill, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if fill == WHITE:
        shape.line.color.rgb = BORDER
        shape.line.width = Pt(1.25)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    if radius:
        try:
            shape.adjustments[0] = 0.06
        except Exception:
            pass
    return shape


def build_pptx(path: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # Header
    _tb(s, 0.6, 0.42, 8.0, 0.6, "16. Team Contribution", size=30, bold=True, color=DARK)
    _tb(s, 0.6, 1.02, 10.0, 0.4,
        "Team of 3 — Shree Deore  ·  Krish Dalvi  ·  Vaishnavi Rathi",
        size=16, color=GRAY)

    # Left: dark card — owned modules
    _card(s, 0.6, 1.7, 5.9, 4.9, fill=NAVY)
    _tb(s, 0.95, 2.0, 5.2, 0.4, "Owned modules", size=17, bold=True, color=PERIWINKLE)
    modules = [
        "Claim verification & fact-checking chatbot",
        "Retrieval & evidence extraction",
        "Supporting / contradicting / neutral verification",
        "Grounded LLM explanation",
        "Complete RAG pipeline integration",
        "Backend/API, security & automated testing",
    ]
    y = 2.5
    for m in modules:
        _tb(s, 0.95, y, 5.2, 0.55, f"•  {m}", size=14.5, color=WHITE)
        y += 0.62

    # Right: light card — how this fits the team
    _card(s, 6.85, 1.7, 5.9, 4.9, fill=WHITE)
    _tb(s, 7.2, 2.0, 5.2, 0.4, "How this fits the team", size=17, bold=True, color=BLUE)
    _tb(s, 7.2, 2.5, 5.2, 1.5,
        "Built as a team of three, we connected the intelligence layer to the "
        "final user experience — the chatbot retrieves evidence, reasons over "
        "it, and explains the result.",
        size=14, color=GRAY)
    # Pill badge
    pill = _card(s, 7.2, 4.35, 3.6, 0.5, fill=LIGHT_BLUE)
    _tb(s, 7.2, 4.44, 3.6, 0.35, "RAG + EVIDENCE + EXPLANATION + API",
        size=12, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
    _tb(s, 7.2, 5.1, 5.2, 1.0,
        "Result: a modular, fully tested pipeline — developed, integrated and "
        "presented by the whole team.",
        size=14, color=GRAY)

    # Footer
    _tb(s, 0.6, 6.95, 6.0, 0.3, "AI-POWERED FAKE NEWS DETECTOR",
        size=10, bold=True, color=GRAY)
    _tb(s, 9.0, 6.95, 3.75, 0.3, "Hackathon Presentation   16",
        size=10, color=GRAY, align=PP_ALIGN.RIGHT)

    prs.save(path)


# ---------------------------------------------------------------------------
# PNG version (1920x1080) — for pasting into the existing deck
# ---------------------------------------------------------------------------

from PIL import Image, ImageDraw, ImageFont  # noqa: E402

W, H = 1920, 1080


def _font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    candidates = [
        r"C:\Windows\Fonts\calibrib.ttf" if bold else r"C:\Windows\Fonts\calibri.ttf",
        r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
    ]
    for path in candidates:
        if os.path.exists(path):
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def _wrap(draw: ImageDraw.ImageDraw, text: str, font, max_w: int) -> list[str]:
    lines: list[str] = []
    cur = ""
    for word in text.split():
        test = f"{cur} {word}".strip()
        if draw.textlength(test, font=font) <= max_w:
            cur = test
        else:
            if cur:
                lines.append(cur)
            cur = word
    if cur:
        lines.append(cur)
    return lines


def build_png(path: str) -> None:
    img = Image.new("RGB", (W, H), "#FFFFFF")
    d = ImageDraw.Draw(img)

    # Header
    d.text((86, 52), "16. Team Contribution", font=_font(46, bold=True), fill="#111827")
    d.text((86, 122), "Team of 3 — Shree Deore  ·  Krish Dalvi  ·  Vaishnavi Rathi",
           font=_font(25), fill="#4B5563")

    # Left card
    d.rounded_rectangle([86, 205, 930, 915], radius=26, fill="#111827")
    d.text((130, 242), "Owned modules", font=_font(27, bold=True), fill="#A5B4FC")
    modules = [
        "Claim verification & fact-checking chatbot",
        "Retrieval & evidence extraction",
        "Supporting / contradicting / neutral verification",
        "Grounded LLM explanation",
        "Complete RAG pipeline integration",
        "Backend/API, security & automated testing",
    ]
    my = 305
    for m in modules:
        d.text((130, my), f"•  {m}", font=_font(23), fill="#FFFFFF")
        my += 92

    # Right card
    d.rounded_rectangle([986, 205, 1830, 915], radius=26, outline="#E5E7EB",
                        fill="#FFFFFF", width=3)
    d.text((1030, 242), "How this fits the team", font=_font(27, bold=True), fill="#2563EB")
    para = ("Built as a team of three, we connected the intelligence layer to the "
            "final user experience — the chatbot retrieves evidence, reasons over "
            "it, and explains the result.")
    para_lines = _wrap(d, para, _font(23), 730)
    py = 310
    for line in para_lines:
        d.text((1030, py), line, font=_font(23), fill="#4B5563")
        py += 38
    # Pill badge
    pill_text = "RAG + EVIDENCE + EXPLANATION + API"
    pill_font = _font(20, bold=True)
    pill_w = int(d.textlength(pill_text, font=pill_font)) + 60
    d.rounded_rectangle([1030, 640, 1030 + pill_w, 700], radius=30, fill="#DBEAFE")
    d.text((1030 + 30, 658), pill_text, font=pill_font, fill="#1D4ED8")
    # Result paragraph
    res = ("Result: a modular, fully tested pipeline — developed, integrated "
           "and presented by the whole team.")
    ry = 745
    for line in _wrap(d, res, _font(23), 730):
        d.text((1030, ry), line, font=_font(23), fill="#4B5563")
        ry += 38

    # Footer
    d.text((86, 960), "AI-POWERED FAKE NEWS DETECTOR", font=_font(16, bold=True), fill="#6B7280")
    footer_right = "Hackathon Presentation   16"
    fr_w = d.textlength(footer_right, font=_font(16))
    d.text((1830 - fr_w, 960), footer_right, font=_font(16), fill="#6B7280")

    img.save(path, "PNG")


if __name__ == "__main__":
    build_pptx("Team_Contribution_Slide.pptx")
    build_png("Team_Contribution_Slide.png")
    print("Saved Team_Contribution_Slide.pptx and Team_Contribution_Slide.png")
